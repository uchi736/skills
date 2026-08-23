# -*- coding: utf-8 -*-
"""
corp_pptx — 自社ハウススタイルのプレゼン生成ライブラリ(corp-slides スキル同梱)v3

■ テンプレートの解決(push安全設計)
    assets/brand/format-16x9.pptx / format-4x3.pptx   ← 公式テンプレ(git管理外。ロゴ・意匠入り)
    → 存在すればそれを土台にする(本物のロゴ・フッター・扉デザインをマスター継承)
    → 無ければ中立なダミー(自前描画・「LOGO」表記・「© Your Company」)で同じ10型を生成
    ブランド資産をリポジトリに含めなくても、コードだけで動作する。

■ 判型
    Deck()                # 16:9(既定)
    Deck(aspect="4:3")    # 4:3
    座標はすべて W/H 比率で計算されるため、両判型で同じAPIが使える。

使い方(最小):
    import sys; sys.path.insert(0, os.path.expanduser("~/.claude/skills/corp-slides/scripts"))
    from corp_pptx import Deck
    d = Deck()
    d.cover("メインタイトル", subtitle="サブ", label="決算説明会", date="2026年7月", dept="○○部")
    d.toc(["背景", "提案", "効果"], current=1)
    d.section(1, "背景")
    d.content("タイトル", lead="結論。", bullets=["要点1", ("親", ["子", ("子2", ["孫"])])])
    d.message("締めのメッセージ")
    d.save("out.pptx")

検証は export_preview.py(PowerPoint COM)で実画像を出して行うこと。
"""
import os, copy, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ============================================================ パレット(自社公開IR資料から実測)
NAVY   = "#01357A"  # コーポレート濃紺(最強の強調)
BLUE   = "#004EA2"  # メインブルー(箱ヘッダ・KPI値・チャート主系列)
BLUE_TX= "#00549A"  # 見出し系テキスト
STEEL  = "#3377AE"  # スチールブルー(目次番号・第2階層)
CYAN   = "#00BCEB"  # 明シアン(アクセントバー・矢印)
CYANSUB= "#66C5DC"  # サブ見出し帯
TEAL   = "#00BB99"  # ティール緑(ポジティブ数値・少量)
RED    = "#C00000"  # ネガティブ数値(△)専用。装飾に使わない
BG     = "#F0F5F8"  # 淡背景ボックス
LTBLUE = "#D6ECF7"  # 選択ハイライト(目次)
PERI   = "#C5D0E6"  # 帯
INK    = "#262626"  # 本文(純黒は使わない)
INK2   = "#4D4D4D"  # 補足・注釈
GREY   = "#A6A6A6"  # 罫線・弱い文字
GRID   = "#E3E7EC"  # チャートのグリッド線
YU     = "Yu Gothic"

SERIES_COLORS = [BLUE, CYANSUB, TEAL, STEEL, NAVY, "#66A6EA", PERI, CYAN]
RAMP = ["#66A6EA", "#2070C7", BLUE, NAVY]

_NUMERIC_RE = re.compile(r"^[+\-△▲]?[\d,.]+\s*(%|円|億円|百万円|人|件|pt)?$")

_ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "assets")
BRAND_FILES = {"16:9": os.path.join(_ASSETS, "brand", "format-16x9.pptx"),
               "4:3":  os.path.join(_ASSETS, "brand", "format-4x3.pptx")}
DUMMY_SIZES = {"16:9": (13.333, 7.5), "4:3": (10.0, 7.5)}

# ダミーモードの中立表記(push されるのはこの文字列だけ)
DUMMY_LOGO   = "LOGO"
DUMMY_FOOTER = "© Your Company All Rights Reserved."
BRAND_CFG_FILE = os.path.join(_ASSETS, "brand", "brand.json")


def load_brand_cfg():
    """assets/brand/brand.json(git管理外)があれば自社固有値を読む。無ければ中立既定。

    brand.json 例: {"company": "株式会社 ○○", "footer": "© ○○ ...", "label": "○○"}
    """
    import json
    try:
        with open(BRAND_CFG_FILE, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}

# 公式レイアウト名 → 用途(ブランドテンプレ側のレイアウト名と一致させること)
LAYOUT_COVER   = "表紙"
LAYOUT_TOC     = "目次"
LAYOUT_SECTION = "扉"
LAYOUT_BODY    = "本文ページ"
LAYOUT_LAST    = "最終ページ"


def _rgb(h):
    return RGBColor.from_string(h.lstrip("#"))

_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


def _set_font(run, size, bold, color, name=YU):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    run.font.color.rgb = _rgb(color)
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def _patch_theme(prs):
    """テーマの和文フォントを游ゴシックへ(ブランド:MS Pゴシック置換 / ダミー:Calibri置換)"""
    for part in prs.part.package.iter_parts():
        pn = str(part.partname)
        if "theme" in pn and pn.endswith(".xml"):
            try:
                blob = part.blob.decode("utf-8")
                blob = blob.replace('typeface="ＭＳ Ｐゴシック"', f'typeface="{YU}"')
                blob = blob.replace('typeface="Calibri Light"', f'typeface="{YU}"')
                blob = blob.replace('typeface="Calibri"', f'typeface="{YU}"')
                blob = blob.replace('<a:ea typeface=""/>', f'<a:ea typeface="{YU}"/>')
                part._blob = blob.encode("utf-8")
            except Exception:
                pass


class Deck:
    """ハウススタイルのデッキビルダー。1メソッド=1スライド(型)。

    aspect : "16:9"(既定) | "4:3"
    brand  : True(既定)=assets/brand のテンプレがあれば使う / False=強制ダミー
    """

    def __init__(self, aspect="16:9", brand=True, template=None):
        if aspect not in DUMMY_SIZES:
            raise ValueError('aspect は "16:9" か "4:3"')
        self.aspect = aspect
        path = template or (BRAND_FILES[aspect] if brand else None)
        self._brand = bool(path and os.path.exists(path))
        if self._brand:
            self.prs = Presentation(path)
            self._clear_slides()
            self._layouts = {lay.name: lay for m in self.prs.slide_masters
                             for lay in m.slide_layouts}
        else:
            self.prs = Presentation()
            w, h = DUMMY_SIZES[aspect]
            self.prs.slide_width = Inches(w)
            self.prs.slide_height = Inches(h)
            self._blank = self.prs.slide_layouts[6]
        _patch_theme(self.prs)
        self.brand_cfg = load_brand_cfg()
        # ---- 判型に依存しない相対ジオメトリ(基準: 公式16:9 = 14.347x8.071)
        self.W = self.prs.slide_width / 914400
        self.H = self.prs.slide_height / 914400
        self.MX = 0.87 if self._brand else 0.061 * self.W
        self.CW = self.W - 2 * self.MX
        self.BODY_TOP = 0.161 * self.H
        self.BODY_BOT = self.H - 0.62
        self._page = 0

    # ---------------------------------------------------------------- 基盤
    def _clear_slides(self):
        sldIdLst = self.prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            rId = sldId.get(qn("r:id"))
            self.prs.part.drop_rel(rId)
            sldIdLst.remove(sldId)

    def _slide(self, layout_name, title=None):
        self._page += 1
        if self._brand:
            lay = self._layouts[layout_name]
            s = self.prs.slides.add_slide(lay)
            # python-pptx はフッター/ページ番号プレースホルダを複製しないため手動継承
            for ph in lay.placeholders:
                if ph.placeholder_format.type in (PP_PLACEHOLDER.FOOTER,
                                                  PP_PLACEHOLDER.SLIDE_NUMBER):
                    s.shapes._spTree.append(copy.deepcopy(ph._element))
            if title is not None and s.shapes.title is not None:
                s.shapes.title.text_frame.text = title
            return s
        # ---- ダミーモード: 中立の意匠を自前で描く
        s = self.prs.slides.add_slide(self._blank)
        W, H = self.W, self.H
        if layout_name == LAYOUT_COVER:
            self.rect(s, 0, 0.510 * H, W, 0.062 * H, PERI)
            self.rect(s, 0, 0.510 * H, W, 0.006 * H, NAVY)
            self.text(s, W - 2.4, 0.510 * H, 1.9, 0.062 * H,
                      [[(DUMMY_LOGO, 22, True, GREY)]], align="right", anchor="middle")
            if title:
                self.text(s, self.MX, 0.365 * H, W - 2 * self.MX, 0.12 * H,
                          [[(title, 28, True, NAVY)]])
        elif layout_name in (LAYOUT_TOC, LAYOUT_SECTION):
            self.rect(s, 0, 0.883 * H, W, 0.045 * H, PERI)
            self.text(s, W - 2.2, 0.870 * H, 1.7, 0.06 * H,
                      [[(DUMMY_LOGO, 18, True, GREY)]], align="right", anchor="middle")
            if title:
                y = 0.065 * H if layout_name == LAYOUT_TOC else 0.400 * H
                size = 24 if layout_name == LAYOUT_TOC else 28
                self.text(s, self.MX, y, W - 2 * self.MX, 0.1 * H,
                          [[(title, size, True, INK)]])
        elif layout_name == LAYOUT_LAST:
            self.text(s, 0, 0.44 * H, W, 0.12 * H,
                      [[(DUMMY_LOGO, 40, True, GREY)]], align="center")
        else:  # 本文ページ
            if title:
                self.text(s, self.MX, 0.045 * H, W - 2 * self.MX - 1.3, 0.08 * H,
                          [[(title, 22, True, INK)]])
            self.rect(s, self.MX, 0.128 * H, self.CW, 0.004 * H, BLUE)
            self.text(s, W - self.MX - 1.3, 0.032 * H, 1.3, 0.06 * H,
                      [[(DUMMY_LOGO, 16, True, GREY)]], align="right")
        # フッター(表紙はページ番号なし)
        self.text(s, self.MX, H - 0.36, 0.5 * W, 0.3,
                  [[(DUMMY_FOOTER, 8, False, GREY)]])
        if layout_name != LAYOUT_COVER:
            self.text(s, W - self.MX - 0.8, H - 0.38, 0.8, 0.3,
                      [[(str(self._page), 10, False, INK2)]], align="right")
        return s

    def rect(self, s, x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE,
             adj=None, dash=False):
        sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.shadow.inherit = False
        if adj is not None:
            try:
                sp.adjustments[0] = adj
            except Exception:
                pass
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid()
            sp.fill.fore_color.rgb = _rgb(fill)
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = _rgb(line)
            sp.line.width = Pt(lw)
            if dash:
                from pptx.enum.dml import MSO_LINE
                sp.line.dash_style = MSO_LINE.DASH
        return sp

    def _tag(self, s, x, y, text, size=11, cx=None):
        """タグ(角丸チップ)。ボックス内のサブ見出しはこの形で示す(全幅の帯にしない — ハウス好み)。
        cx を渡すと中央揃えで配置。戻り値は描画幅。"""
        t = str(text)
        tw = 0.30 + sum(0.16 if ord(c) > 0x2E80 else 0.088 for c in t) * (size / 11)
        if cx is not None:
            x = cx - tw / 2
        self.rect(s, x, y, tw, 0.34, "#FFFFFF", line=GREY, lw=1.0,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.5)
        self.text(s, x, y, tw, 0.34, [[(t, size, False, INK2)]],
                  align="center", anchor="middle")
        return tw

    def _line_seg(self, s, x1, y1, x2, y2, color, wpt=1.5):
        """任意角度の線分(細長い矩形を回転させて描く)。アイコン・概念図用。"""
        import math
        L = max(math.hypot(x2 - x1, y2 - y1), 0.01)
        th = wpt / 72.0
        sp = self.rect(s, (x1 + x2) / 2 - L / 2, (y1 + y2) / 2 - th / 2, L, th, color)
        sp.rotation = math.degrees(math.atan2(y2 - y1, x2 - x1))
        return sp

    def icon(self, s, kind, x, y, size=0.5, color=BLUE):
        """簡易アイコン(単色・幾何図形のみ。多色・立体・絵文字は使わない)。
        kind: book(辞書・規程) | cycle(学習・循環) | network(グラフ・関係) |
              search(検索) | db(データ) | doc(文書) | gear(処理)
        """
        if kind == "book":
            self.rect(s, x, y + size * 0.10, size, size * 0.78, color,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.15)
            self.rect(s, x + size * 0.47, y + size * 0.16, size * 0.06, size * 0.66, "#FFFFFF")
            for i in range(2):
                yy = y + size * (0.34 + 0.20 * i)
                self.rect(s, x + size * 0.13, yy, size * 0.24, 0.016, "#FFFFFF")
                self.rect(s, x + size * 0.63, yy, size * 0.24, 0.016, "#FFFFFF")
        elif kind == "cycle":
            self.rect(s, x, y, size, size, color, shape=MSO_SHAPE.CIRCULAR_ARROW)
        elif kind == "network":
            pts = [(0.50, 0.12), (0.12, 0.58), (0.88, 0.58), (0.50, 0.88)]
            for a, b in [(0, 1), (0, 2), (1, 2), (1, 3), (2, 3)]:
                self._line_seg(s, x + pts[a][0] * size, y + pts[a][1] * size,
                               x + pts[b][0] * size, y + pts[b][1] * size, color, 1.6)
            d0 = size * 0.30
            for px, py in pts:
                self.rect(s, x + px * size - d0 / 2, y + py * size - d0 / 2, d0, d0,
                          color, shape=MSO_SHAPE.OVAL)
        elif kind == "search":
            d0 = size * 0.66
            self.rect(s, x, y, d0, d0, None, line=color, lw=3.0, shape=MSO_SHAPE.OVAL)
            self._line_seg(s, x + d0 * 0.85, y + d0 * 0.85, x + size, y + size, color, 3.0)
        elif kind == "db":
            self.rect(s, x + size * 0.12, y, size * 0.76, size, color, shape=MSO_SHAPE.CAN)
        elif kind == "doc":
            self.rect(s, x + size * 0.12, y, size * 0.76, size, color,
                      shape=MSO_SHAPE.FOLDED_CORNER)
            for i in range(3):
                yy = y + size * (0.30 + 0.18 * i)
                self.rect(s, x + size * 0.26, yy, size * 0.40, 0.016, "#FFFFFF")
        else:
            self.rect(s, x, y, size, size, color, shape=MSO_SHAPE.GEAR_6)

    # ---------------------------------------------------------------- 概念図の部品(全て編集可能なネイティブ図形)
    def card(self, s, x, y, w, h, fill="#FFFFFF", line=GRID, lw=1.0, dash=False, adj=0.06):
        """角丸カード(ヘアライン枠)。概念図の基本コンテナ。dash=True で計画・棄却の表現。"""
        return self.rect(s, x, y, w, h, fill, line=line, lw=lw,
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=adj, dash=dash)

    def zone(self, s, x, y, w, h, title):
        """見出し付きゾーン枠(構築/検索のような対構造に)。中身は低レベルAPIで組む。"""
        self.card(s, x, y, w, h, fill=BG)
        self.text(s, x + 0.22, y + 0.16, w - 0.44, 0.32, [[(title, 11, True, BLUE_TX)]])

    def tri_down(self, s, cx, y, w=0.30, h=0.26, color=CYAN):
        sp = self.rect(s, cx - w / 2, y, w, h, color, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
        sp.rotation = 180
        return sp

    def tag_row(self, s, x, y, texts, size=9.5):
        """タグを横に並べる。戻り値は右端x。"""
        for t in texts:
            x += self._tag(s, x, y, t, size=size) + 0.15
        return x

    def node(self, s, cx, cy, label, r=0.34, fill=BLUE, dash=False, size=9.5):
        """ノードリンク図の丸ノード。dash=True は「欠けている要素・未接続」の表現。"""
        if dash:
            self.rect(s, cx - r, cy - r, r * 2, r * 2, "#FFFFFF", line=GREY, lw=1.2,
                      shape=MSO_SHAPE.OVAL, dash=True)
            col = INK2
        else:
            self.rect(s, cx - r, cy - r, r * 2, r * 2, fill, shape=MSO_SHAPE.OVAL)
            col = "#FFFFFF"
        self.text(s, cx - r, cy - r, r * 2, r * 2,
                  [[(ln, size, True, col)] for ln in str(label).split("\n")],
                  align="center", anchor="middle", space_after=0, line_spacing=1.1)

    def edge(self, s, x1, y1, x2, y2, label=None, dash=False, color=STEEL):
        """ノード間の線(先にedge、後からnodeを描くと重なりが綺麗)。dash=弱い/欠けた関係。"""
        import math
        if dash:
            L = math.hypot(x2 - x1, y2 - y1)
            nseg = max(int(L / 0.09), 2)
            for k in range(0, nseg, 2):
                t0, t1 = k / nseg, min((k + 1) / nseg, 1.0)
                self._line_seg(s, x1 + (x2 - x1) * t0, y1 + (y2 - y1) * t0,
                               x1 + (x2 - x1) * t1, y1 + (y2 - y1) * t1, GREY, 1.5)
        else:
            self._line_seg(s, x1, y1, x2, y2, color, 1.5)
        if label:
            self.text(s, (x1 + x2) / 2 - 0.9, (y1 + y2) / 2 - 0.24, 1.8, 0.22,
                      [[(label, 8, False, INK2)]], align="center")

    def step_flow(self, s, x, y, w, h, steps, intro=None, last_hl=True):
        """番号付き処理フロー(①→②→…)。手法の動作説明・アプローチ図解の上段の定番。

        steps: [(タイトル, 説明), ...] 3〜5個
        intro: (小ラベル, 本文) — 先頭に置く入力カード(質問の実例など)。省略可
        last_hl: 最終ステップを淡青で強調(成果の位置)
        """
        aw = 0.34
        n = len(steps)
        iw = 2.1 if intro else 0
        cw_ = (w - iw - (n - 1 + (1 if intro else 0)) * aw) / n
        cx = x
        if intro:
            self.card(s, cx, y, iw, h, line=STEEL, lw=1.5)
            self.text(s, cx + 0.18, y + 0.14, iw - 0.36, 0.28, [[(intro[0], 9, False, INK2)]])
            self.text(s, cx + 0.18, y + 0.46, iw - 0.36, h - 0.6,
                      [[(intro[1], 12, True, INK)]], line_spacing=1.3)
            cx += iw
            self.tri_right(s, cx + 0.08, y + h / 2 - 0.11, aw - 0.16, 0.22, CYAN)
            cx += aw
        for i, (t, desc) in enumerate(steps):
            self.card(s, cx, y, cw_, h,
                      fill=(LTBLUE if (last_hl and i == n - 1) else BG))
            self.rect(s, cx + 0.16, y + 0.15, 0.28, 0.28, BLUE, shape=MSO_SHAPE.OVAL)
            self.text(s, cx + 0.16, y + 0.15, 0.28, 0.28,
                      [[(str(i + 1), 10, True, "#FFFFFF")]], align="center", anchor="middle")
            self.text(s, cx + 0.52, y + 0.14, cw_ - 0.66, 0.30, [[(t, 11.5, True, BLUE_TX)]])
            self.text(s, cx + 0.16, y + 0.54, cw_ - 0.32, h - 0.68,
                      [[(desc, 9, False, INK2)]], line_spacing=1.25)
            cx += cw_
            if i < n - 1:
                self.tri_right(s, cx + 0.08, y + h / 2 - 0.11, aw - 0.16, 0.22, CYAN)
                cx += aw

    def stage_flow(self, s, x, y, w, h, stages, label=None):
        """状態付きステージパイプライン(学習手順・移行計画などの段階図)。

        stages: [(タイトル, サブ(改行は\\n), state), ...]
          state: "done"=青 | "now"=スチール青(現行) | "dead"=白+灰破線+×(棄却) | "plan"=白+青破線(計画)
        dead ステージから出る矢印には「切替」ラベルが自動で付く。
        """
        if label:
            self.text(s, x, y - 0.30, w, 0.26, [[(label, 9.5, False, INK2)]])
        aw = 0.36
        n = len(stages)
        cw_ = (w - (n - 1) * aw) / n
        cx = x
        for i, (t, sub, state) in enumerate(stages):
            if state == "dead":
                self.card(s, cx, y, cw_, h, line=GREY, lw=1.2, dash=True)
                self.text(s, cx, y + 0.15, cw_, 0.32,
                          [[(t + " ", 13, True, INK2), ("×", 15, True, GREY)]], align="center")
                tcol = INK2
            elif state == "plan":
                self.card(s, cx, y, cw_, h, line=STEEL, lw=1.2, dash=True)
                self.text(s, cx, y + 0.15, cw_, 0.32, [[(t, 12, True, BLUE_TX)]], align="center")
                tcol = BLUE_TX
            else:
                self.card(s, cx, y, cw_, h, fill=(STEEL if state == "now" else BLUE), line=None)
                self.text(s, cx, y + 0.15, cw_, 0.32, [[(t, 13, True, "#FFFFFF")]], align="center")
                tcol = "#FFFFFF"
            if sub:
                self.text(s, cx + 0.10, y + 0.52, cw_ - 0.20, h - 0.62,
                          [[(ln, 9, False, tcol)] for ln in str(sub).split("\n")],
                          align="center", space_after=2, line_spacing=1.25)
            cx += cw_
            if i < n - 1:
                if state == "dead":
                    self.text(s, cx - 0.30, y + h / 2 - 0.44, aw + 0.60, 0.24,
                              [[("切替", 9.5, True, BLUE_TX)]], align="center")
                self.tri_right(s, cx + 0.08, y + h / 2 - 0.11, aw - 0.16, 0.22,
                               GREY if stages[i + 1][2] == "plan" else CYAN)
                cx += aw

    def figure(self, title, image, lead=None, source=None):
        """⑭ 図版スライド。HTML図解(scripts/html2png.py でPNG化)などの画像を
        本文エリア中央に比率維持で配置する。
        編集可能なマスター(.html)は画像と同じ場所に必ず残すこと(画像は直接編集不可)。"""
        from PIL import Image as _PILImage
        s = self._slide(LAYOUT_BODY, title=title)
        y0 = self._lead(s, lead)
        iw, ih = _PILImage.open(image).size
        top = y0 + 0.10
        bot = self.BODY_BOT - (0.30 if source else 0.12)
        r = min(self.CW / iw, (bot - top) / ih)
        w, h = iw * r, ih * r
        s.shapes.add_picture(image, Inches(self.MX + (self.CW - w) / 2),
                             Inches(top + (bot - top - h) / 2), Inches(w), Inches(h))
        self._source(s, source)
        return s

    def tri_right(self, s, x, y, w, h, fill):
        """右向き三角形(▶)。図解の「流れ」はブロック矢印ではなくこれで示す(ハウス好み)。
        x,y,w,h は見た目のバウンディングボックス(回転を吸収して配置する)。"""
        cx, cy = x + w / 2, y + h / 2
        sp = self.rect(s, cx - h / 2, cy - w / 2, h, w, fill,
                       shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
        sp.rotation = 90
        return sp

    def text(self, s, x, y, w, h, paras, align="left", anchor="top",
             space_after=8, line_spacing=1.15, wrap=True):
        """paras: [[(text,size,bold,color[,font]), ...], ...]"""
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = _ANCHOR[anchor]
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, para in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = _ALIGN[align]
            p.space_after = Pt(space_after)
            p.line_spacing = line_spacing
            for run_spec in para:
                txt, size, bold, color = run_spec[:4]
                name = run_spec[4] if len(run_spec) > 4 else YU
                r = p.add_run()
                r.text = txt
                _set_font(r, size, bold, color, name)
        return tb

    def _lead(self, s, lead, y=None):
        """キーメッセージ(リード文): ■ + 18〜20pt。短文は20pt、2行になる長文は18pt。"""
        y = self.BODY_TOP if y is None else y
        if not lead:
            return y + 0.85
        size = 20 if len(lead) <= 40 else 18
        self.text(s, self.MX + 0.05, y + 0.10, self.CW - 0.15, 1.1,
                  [[("■", int(size * 0.8), True, INK), ("  " + lead, size, False, INK)]])
        chars_per_line = int(self.CW * 60 / size)   # 全角換算の目安
        lines = max(1, -(-len(lead) // chars_per_line))
        line_h = size / 72 * 1.35
        return y + 0.10 + line_h * lines + 0.45

    # ---------------------------------------------------------------- 階層マーカー(唯一の定義)
    # 全型共通の3段ラダー: 本文 ●18pt → 　□15pt → 　　‐13pt / 記号は黒(INK)で控えめに
    # 12pt → 11pt → 10pt(基準時。縮小文脈では比率維持)。他のマーカーは使わない。
    def bullet(self, text, level=1, size=18, bold=False):
        if level <= 1:
            return [("●", max(8, round(size * 12 / 18)), True, INK), ("  " + text, size, bold, INK)]
        if level == 2:
            s = max(10, size - 3)
            return [("　□", max(8, round(s * 11 / 15)), True, INK), ("  " + text, s, bold, INK)]
        s = max(9, size - 5)
        return [("　　‐", max(8, round(s * 10 / 13)), False, INK), ("  " + text, s, bold, INK2)]

    def nest_paras(self, items, size=18, level=1):
        """ネストしたリストを段落列に変換。items: ["文", ("親", ["子", ("子2", ["孫"])]), ...]"""
        paras = []
        for it in items:
            if (isinstance(it, (tuple, list)) and len(it) == 2
                    and isinstance(it[1], (list, tuple))):
                head, kids = it
                paras.append(self.bullet(head, level, size))
                paras.extend(self.nest_paras(kids, size, min(level + 1, 3)))
            else:
                paras.append(self.bullet(it, level, size))
        return paras

    def bullet_para(self, t, size=15, color=INK, bold=False):
        """旧API互換。新規コードは bullet()/nest_paras() を使う"""
        return self.bullet(t, 1, size, bold)

    def note(self, text):
        s = self.prs.slides[-1]
        s.notes_slide.notes_text_frame.text = text

    def _source(self, s, source):
        if source:
            self.text(s, self.MX, self.H - 0.78, self.CW, 0.3,
                      [[("出典:" + source, 10, False, INK2)]])

    # ---------------------------------------------------------------- ① 表紙
    def cover(self, title, subtitle=None, label=None, date=None, dept=None, company=None):
        s = self._slide(LAYOUT_COVER, title=title)
        H = self.H
        if company is None:
            company = self.brand_cfg.get("company", "[会社名]")
        if label:
            self.text(s, self.MX, 0.316 * H, 0.8 * self.W, 0.4, [[(label, 15, False, INK)]])
        if subtitle:
            self.text(s, self.MX, 0.558 * H, 0.85 * self.W, 0.5,
                      [[("― " + subtitle + " ―", 17, True, BLUE_TX)]])
        if date:
            self.text(s, self.MX, 0.651 * H, 7.0, 0.4, [[(date, 14, False, INK)]])
        self.text(s, self.MX, 0.768 * H, 9.0, 0.4, [[(company, 16, True, INK)]])
        if dept:
            self.text(s, self.MX, 0.820 * H, 9.0, 0.4, [[(dept, 13, False, INK)]])
        return s

    # ---------------------------------------------------------------- ② 目次
    def toc(self, items, current=None, title="目次"):
        s = self._slide(LAYOUT_TOC, title=title)
        x = 0.237 * self.W
        y = 0.242 * self.H
        step = 0.118 * self.H
        for i, it in enumerate(items, 1):
            name, sub = (it if isinstance(it, (tuple, list)) else (it, None))
            if current == i:
                hl_w = min(self.W - x - 0.9,
                           1.3 + len(name) * 0.34 + (len(sub) * 0.20 if sub else 0) + 0.8)
                self.rect(s, x - 0.28, y - 0.14, hl_w, 0.72, LTBLUE)
            runs = [(f"{i}.", 21, True, STEEL), ("   " + name, 23, True, BLUE_TX)]
            if sub:
                runs.append(("  ― " + sub, 14, True, STEEL))
            self.text(s, x, y, self.W - x - 0.9, 0.5, [runs])
            y += step
        return s

    # ---------------------------------------------------------------- ③ 中扉
    def section(self, number, title):
        return self._slide(LAYOUT_SECTION, title=f"{number}. {title}")

    # ---------------------------------------------------------------- ④ 本文・箇条書き
    def content(self, title, lead=None, bullets=None, source=None):
        """bullets: ["文"] または [("親", ["子1", ("子2", ["孫"])])] — 最大3階層"""
        s = self._slide(LAYOUT_BODY, title=title)
        y = self._lead(s, lead)
        paras = self.nest_paras(bullets or [], size=18)
        if paras:
            self.text(s, self.MX + 0.15, y, self.CW - 0.3, self.BODY_BOT - y,
                      paras, space_after=12)
        self._source(s, source)
        return s

    # ---------------------------------------------------------------- ⑤ ボックス対比(2〜3列)
    def boxes(self, title, boxes, lead=None, arrow=False):
        s = self._slide(LAYOUT_BODY, title=title)
        y0 = self._lead(s, lead) - 0.15
        n = len(boxes)
        gap = 0.58 if arrow else 0.30
        w = (self.CW - gap * (n - 1)) / n
        bot = self.BODY_BOT - 0.10
        for i, (head, sub, items) in enumerate(boxes):
            x = self.MX + i * (w + gap)
            self.rect(s, x, y0, w, 0.55, BLUE)
            self.text(s, x, y0, w, 0.55, [[(head, 17, True, "#FFFFFF")]],
                      align="center", anchor="middle")
            yy = y0 + 0.55
            self.rect(s, x, yy, w, bot - yy, BG)  # 本文はヘッダに密着(1枚のカードに見せる)
            iy = yy + 0.18
            if sub:
                self._tag(s, x + 0.30, iy, sub)
                iy += 0.48
            self.text(s, x + 0.30, iy, w - 0.60, bot - iy - 0.17,
                      self.nest_paras(items, size=12), space_after=8)
            if arrow and i < n - 1:
                self.tri_right(s, x + w + 0.08, (y0 + bot) / 2 - 0.26,
                               gap - 0.16, 0.52, CYAN)
        return s

    # ---------------------------------------------------------------- ⑥ KPIハイライト
    def kpi(self, title, kpis, lead=None, source=None):
        s = self._slide(LAYOUT_BODY, title=title)
        y0 = self._lead(s, lead) - 0.05
        n = len(kpis)
        gap = 0.32
        w = (self.CW - gap * (n - 1)) / n
        h = min(3.60, self.BODY_BOT - 0.35 - y0)
        for i, item in enumerate(kpis):
            label, value, unit, note = (list(item) + [None] * 4)[:4]
            x = self.MX + i * (w + gap)
            self.rect(s, x, y0, w, h, BG)
            self.rect(s, x, y0, w, 0.07, BLUE)
            self.text(s, x + 0.2, y0 + 0.30, w - 0.4, 0.5,
                      [[(label, 14, True, INK)]], align="center")
            vcolor = BLUE
            if str(value).startswith(("+", "▲")):
                vcolor = TEAL
            elif str(value).startswith(("-", "−", "△", "▼")):
                vcolor = RED
            vsize = 40 if len(str(value)) <= 6 else (32 if len(str(value)) <= 9 else 26)
            runs = [(str(value), vsize, True, vcolor)]
            if unit:
                runs.append((" " + unit, 16, True, INK2))
            self.text(s, x + 0.2, y0 + h / 2 - 0.42, w - 0.4, 1.0, [runs],
                      align="center", anchor="middle")
            if note:
                self.text(s, x + 0.2, y0 + h - 0.62, w - 0.4, 0.5,
                          [[(note, 11, False, INK2)]], align="center")
        self._source(s, source)
        return s

    # ---------------------------------------------------------------- ⑦ ネイティブチャート
    def chart(self, title, kind, categories, series, lead=None, unit=None,
              points=None, value_fmt="#,##0", source=None):
        s = self._slide(LAYOUT_BODY, title=title)
        y0 = self._lead(s, lead)
        cw = self.CW * 0.64 if points else self.CW
        if unit:
            self.text(s, self.MX, y0 - 0.05, 3.0, 0.3,
                      [[(f"(単位:{unit})", 10, False, INK2)]])
        ch_y = y0 + 0.25
        ch_h = self.BODY_BOT - 0.10 - ch_y
        if kind == "waterfall":
            self._waterfall(s, self.MX, ch_y, cw, ch_h, categories, series[0][1])
            self._points_panel(s, ch_y, ch_h, cw, points)
            self._source(s, source)
            return s
        data = CategoryChartData()
        data.categories = categories
        for name, vals in series:
            data.add_series(name, vals)
        xl = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
              "stack": XL_CHART_TYPE.COLUMN_STACKED,
              "line": XL_CHART_TYPE.LINE}[kind]
        gf = s.shapes.add_chart(xl, Inches(self.MX), Inches(ch_y), Inches(cw), Inches(ch_h), data)
        c = gf.chart
        c.has_title = False
        c.font.size = Pt(10)
        c.font.name = YU
        c.font.color.rgb = _rgb(INK)
        c.has_legend = len(series) > 1
        if c.has_legend:
            c.legend.position = XL_LEGEND_POSITION.BOTTOM
            c.legend.include_in_layout = False
        plot = c.plots[0]
        try:
            plot.gap_width = 60 if kind != "line" else 150
        except Exception:
            pass
        for i, ser in enumerate(plot.series):
            col = _rgb(SERIES_COLORS[i % len(SERIES_COLORS)])
            if kind == "line":
                ser.format.line.color.rgb = col
                ser.format.line.width = Pt(2.5)
                ser.smooth = False
            else:
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = col
                ser.format.line.fill.background()
        if kind == "bar" and len(series) == 1:
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.number_format = value_fmt
            dl.number_format_is_linked = False
            dl.position = XL_LABEL_POSITION.OUTSIDE_END
            dl.font.size = Pt(10); dl.font.bold = True
            dl.font.color.rgb = _rgb(INK); dl.font.name = YU
        elif kind == "stack":
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.number_format = value_fmt
            dl.number_format_is_linked = False
            dl.font.size = Pt(9); dl.font.bold = True
            dl.font.color.rgb = _rgb("#FFFFFF"); dl.font.name = YU
        cat = c.category_axis
        cat.has_major_gridlines = False
        cat.format.line.color.rgb = _rgb(GREY)
        cat.tick_labels.font.size = Pt(10)
        cat.tick_labels.font.name = YU
        val = c.value_axis
        if kind != "line":
            val.minimum_scale = 0  # 棒グラフの0始まりは省略不可(差の誇張防止)
        val.has_major_gridlines = True
        val.major_gridlines.format.line.color.rgb = _rgb(GRID)
        val.major_gridlines.format.line.width = Pt(0.75)
        val.format.line.fill.background()
        val.tick_labels.font.size = Pt(9)
        val.tick_labels.font.name = YU
        val.tick_labels.number_format = value_fmt
        val.tick_labels.number_format_is_linked = False
        self._points_panel(s, ch_y, ch_h, cw, points)
        self._source(s, source)
        return s

    def _points_panel(self, s, ch_y, ch_h, cw, points):
        """chart 系の右側「ポイント」欄(points が無ければ何もしない)"""
        if not points:
            return
        px = self.MX + cw + 0.32
        pw = self.W - self.MX - px
        self.rect(s, px, ch_y, pw, 0.46, NAVY)
        self.text(s, px, ch_y, pw, 0.46, [[("ポイント", 13, True, "#FFFFFF")]],
                  align="center", anchor="middle")
        self.rect(s, px, ch_y + 0.46, pw, ch_h - 0.46, BG)
        self.text(s, px + 0.22, ch_y + 0.72, pw - 0.44, ch_h - 0.95,
                  self.nest_paras(points, size=12), space_after=10)

    def _waterfall(self, s, x, y, w, h, categories, values):
        """滝グラフ(ネイティブ図形で描画)。増減プロセスの分解(利益ブリッジ等)専用。

        values: [期首実額, ±増減..., 期末実額(None なら自動計算)]。期首・期末は正の実額。
        色は 期首/期末=BLUE、増加=TEAL、減少=RED(負値専用色の本来の用途)。
        全バーにデータラベルが付くため値軸・グリッドは引かない(IR流の全数表示)。
        """
        vals = list(values)
        if vals[-1] is None:
            vals[-1] = vals[0] + sum(vals[1:-1])
        n = len(vals)
        levels, cum = [(0.0, float(vals[0]))], float(vals[0])
        for d in vals[1:-1]:
            b, t = (cum, cum + d) if d >= 0 else (cum + d, cum)
            levels.append((b, t))
            cum += d
        levels.append((0.0, float(vals[-1])))
        ymax = max(t for _, t in levels) * 1.18
        lab_h, cat_h = 0.26, 0.34
        plot_top = y + lab_h
        plot_h = h - lab_h - cat_h
        base_y = plot_top + plot_h

        def py(v):
            return plot_top + plot_h * (1 - v / ymax)

        bw = w / (n + (n - 1) * 0.5)     # ギャップはバー幅の1/2
        step = bw * 1.5
        for i, (b, t) in enumerate(levels):
            bx = x + i * step
            is_total = (i == 0 or i == n - 1)
            d = vals[i]
            fill = BLUE if is_total else (TEAL if d >= 0 else RED)
            self.rect(s, bx, py(t), bw, max(py(b) - py(t), 0.02), fill)
            if is_total:
                lab, lc = f"{vals[i]:,.0f}", INK
            elif d >= 0:
                lab, lc = f"+{d:,.0f}", TEAL
            else:
                lab, lc = f"△{abs(d):,.0f}", RED
            self.text(s, bx - step * 0.25, py(t) - 0.26, bw + step * 0.5, 0.24,
                      [[(lab, 10, True, lc)]], align="center", anchor="bottom")
            self.text(s, bx - step * 0.25, base_y + 0.06, bw + step * 0.5, cat_h - 0.06,
                      [[(str(categories[i]), 10, False, INK)]], align="center")
            if i < n - 1:  # 累積レベルを次の棒へつなぐ細線
                lev = t if (is_total or d >= 0) else b
                self.rect(s, bx + bw, py(lev) - 0.005, step - bw, 0.01, GREY)
        self.rect(s, x, base_y - 0.005, w, 0.01, GREY)

    # ---------------------------------------------------------------- ⑧ ネイティブ表
    def table(self, title, headers, rows, lead=None, col_widths=None,
              first_col_header=True, font_size=12, unit=None, source=None,
              highlight_row=None):
        s = self._slide(LAYOUT_BODY, title=title)
        y0 = self._lead(s, lead)
        if unit:
            self.text(s, self.MX, y0 - 0.05, self.CW, 0.3,
                      [[(f"(単位:{unit})", 10, False, INK2)]], align="right")
        ty = y0 + 0.26
        ncols, nrows = len(headers), len(rows) + 1
        row_h = min(0.48, (self.BODY_BOT - 0.10 - ty) / nrows)
        gf = s.shapes.add_table(nrows, ncols, Inches(self.MX), Inches(ty),
                                Inches(self.CW), Inches(row_h * nrows))
        tbl = gf.table
        tbl.first_row = False
        tbl.horz_banding = False
        if col_widths:
            total = sum(col_widths)
            for i, cwd in enumerate(col_widths):
                tbl.columns[i].width = Emu(int(Inches(self.CW) * cwd / total))
        for r in range(nrows):
            tbl.rows[r].height = Inches(row_h)
        body_size = font_size

        def put(cell, txt, size, bold, color, fill, align):
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(fill if fill else "#FFFFFF")
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = _ALIGN[align]
            r = p.add_run()
            r.text = str(txt)
            _set_font(r, size, bold, color)

        for cji, htxt in enumerate(headers):
            put(tbl.cell(0, cji), htxt, body_size, True, "#FFFFFF", BLUE, "center")
        for ri, row in enumerate(rows, start=1):
            zebra = BG if ri % 2 == 0 else None
            is_hl = (highlight_row is not None and ri - 1 == highlight_row)
            for cji, val in enumerate(row):
                txt = str(val)
                is_num = bool(_NUMERIC_RE.match(txt.replace(",", "").replace(" ", "")))
                color = INK
                bold = is_hl
                fill = LTBLUE if is_hl else zebra
                align = "right" if is_num and cji > 0 else "left"
                if cji == 0 and first_col_header:
                    bold = True
                    fill = LTBLUE if (is_hl or ri % 2 == 0) else "#E8F2FA"
                    align = "left"
                if txt.startswith(("△", "-", "−", "▼")) and is_num:
                    color = RED
                put(tbl.cell(ri, cji), txt, body_size, bold, color, fill, align)
        self._source(s, source)
        return s

    # ---------------------------------------------------------------- ⑪ 評価マトリクス(コンサル型)
    def matrix_eval(self, title, options, criteria, ratings, lead=None,
                    verdicts=None, verdict_header="位置づけ", recommend=None,
                    source=None):
        """選定・評価マトリクス(コンサル比較表)。候補を列に立て、左に比較軸の行を並べる。

        options : 列(候補)の名前リスト(最大4)
        criteria: 行(比較軸)の名前リスト
        ratings : 候補ごとの値リスト(ratings[i][j] = 候補i×軸j)。
                  値は "◎"/"○"/"△"/"×" の記号、短文テキスト、["箇条", ...] のリストを混在可
        verdicts: 最下段の位置づけ・結論行(省略可)
        recommend: 強調する候補列の index(0始まり)
        """
        s = self._slide(LAYOUT_BODY, title=title)
        y0 = self._lead(s, lead)
        ty = y0 + 0.15
        ncols = 1 + len(options)
        nrows = 1 + len(criteria) + (1 if verdicts else 0)
        h_avail = self.BODY_BOT - 0.15 - ty
        head_h = 0.50
        row_h = max(0.42, min(1.05, (h_avail - head_h) / (nrows - 1)))
        gf = s.shapes.add_table(nrows, ncols, Inches(self.MX), Inches(ty),
                                Inches(self.CW), Inches(head_h + row_h * (nrows - 1)))
        tbl = gf.table
        tbl.first_row = False
        tbl.horz_banding = False
        lab_w = 1.55
        opt_w = (self.CW - lab_w) / len(options)
        for i, wd in enumerate([lab_w] + [opt_w] * len(options)):
            tbl.columns[i].width = Emu(int(Inches(wd)))
        tbl.rows[0].height = Inches(head_h)
        for r in range(1, nrows):
            tbl.rows[r].height = Inches(row_h)

        def put(cell, paras, fill):
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(fill if fill else "#FFFFFF")
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            for k, (runs, align) in enumerate(paras):
                p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                p.alignment = _ALIGN[align]
                p.space_after = Pt(2)
                p.line_spacing = 1.1
                for txt, size, bold, color in runs:
                    r = p.add_run()
                    r.text = str(txt)
                    _set_font(r, size, bold, color)

        SYM = {"◎": (BLUE, True), "○": (INK, False), "△": (INK2, False),
               "×": (INK2, False), "－": (INK2, False), "-": (INK2, False)}

        def body_paras(val):
            if isinstance(val, (list, tuple)):
                return [([("●", 6.5, False, INK), ("　" + str(it), 9.5, False, INK)], "left")
                        for it in val]
            v = str(val)
            if v in SYM:
                color, bold = SYM[v]
                return [([(v, 16, bold, color)], "center")]
            return [([(v, 10, False, INK)], "center")]

        # ヘッダ行(候補)。左上コーナーは空白のまま
        put(tbl.cell(0, 0), [([("", 11, True, "#FFFFFF")], "center")], None)
        for j, opt in enumerate(options, start=1):
            is_rec = (recommend is not None and j - 1 == recommend)
            put(tbl.cell(0, j), [([(opt, 12, True, "#FFFFFF")], "center")],
                NAVY if is_rec else BLUE)
        for i, cri in enumerate(criteria):
            ri = i + 1
            zebra = BG if ri % 2 == 0 else None
            put(tbl.cell(ri, 0), [([(cri, 11, True, "#FFFFFF")], "center")], BLUE)
            for j in range(len(options)):
                is_rec = (recommend is not None and j == recommend)
                put(tbl.cell(ri, j + 1), body_paras(ratings[j][i]),
                    LTBLUE if is_rec else zebra)
        if verdicts:
            ri = nrows - 1
            put(tbl.cell(ri, 0), [([(verdict_header, 11, True, "#FFFFFF")], "center")], BLUE)
            for j, v in enumerate(verdicts):
                is_rec = (recommend is not None and j == recommend)
                put(tbl.cell(ri, j + 1),
                    [([(v, 10.5, is_rec, BLUE_TX if is_rec else INK)], "center")],
                    LTBLUE if is_rec else "#E8F2FA")
        self._source(s, source)
        return s

    # ---------------------------------------------------------------- ⑫ 2軸マトリクス(ポジショニング)
    def matrix_2x2(self, title, x_axis, y_axis, items, lead=None,
                   quadrants=None, emphasize=None, source=None):
        """2軸ポジショニングマップ。

        x_axis / y_axis: (低ラベル, 高ラベル, 軸名)
        items    : [(名前, x, y)] x,y は 0.0〜1.0(右上が高・高)
        quadrants: [左下, 右下, 左上, 右上] の象限ラベル(省略可)
        emphasize: 強調する象限 index(0=左下,1=右下,2=左上,3=右上)
        """
        s = self._slide(LAYOUT_BODY, title=title)
        y0 = self._lead(s, lead)
        px = self.MX + 1.05
        py = y0 + 0.30
        pw = self.CW - 1.55
        ph = self.BODY_BOT - 0.60 - py
        # 象限の背景(強調象限は淡青)
        for qi, (qx, qy) in enumerate([(0, 1), (1, 1), (0, 0), (1, 0)]):  # 左下,右下,左上,右上
            fill = LTBLUE if emphasize == qi else BG
            self.rect(s, px + qx * pw / 2, py + qy * ph / 2, pw / 2, ph / 2, fill)
        # 外枠と十字
        self.rect(s, px, py, pw, ph, None, line=GREY, lw=1.0)
        self.rect(s, px + pw / 2 - 0.008, py, 0.016, ph, "#FFFFFF")
        self.rect(s, px, py + ph / 2 - 0.008, pw, 0.016, "#FFFFFF")
        # 象限ラベル(角に小さく)
        if quadrants:
            pos = [(px + 0.15, py + ph - 0.42, "left"), (px + pw / 2 + 0.15, py + ph - 0.42, "left"),
                   (px + 0.15, py + 0.12, "left"), (px + pw / 2 + 0.15, py + 0.12, "left")]
            for qi, q in enumerate(quadrants):
                if q:
                    qx, qy, al = pos[qi]
                    bold = (emphasize == qi)
                    self.text(s, qx, qy, pw / 2 - 0.3, 0.35,
                              [[(q, 11, bold, BLUE_TX if bold else GREY)]], align=al)
        # 軸ラベル
        lo_x, hi_x, name_x = x_axis
        lo_y, hi_y, name_y = y_axis
        self.text(s, px, py + ph + 0.08, 1.2, 0.3, [[(lo_x, 10, False, INK2)]])
        self.text(s, px + pw - 1.2, py + ph + 0.08, 1.2, 0.3,
                  [[(hi_x, 10, False, INK2)]], align="right")
        self.text(s, px, py + ph + 0.08, pw, 0.32,
                  [[(name_x, 12, True, INK)]], align="center")
        # rotation=270 のテキストは下から上に読まれるため、低→軸名→高 の順で組む
        ytb = self.text(s, px - 0.62, py + ph / 2 - 1.2, 0.5, 2.4,
                        [[(lo_y + "  ← " + name_y + " →  " + hi_y, 11, True, INK)]],
                        align="center", anchor="middle", wrap=False)
        ytb.rotation = 270
        # アイテム(チップ)
        for item in items:
            name, ix, iy = item[:3]
            w = 0.34 * len(name) / 2 + 0.75
            h = 0.44
            cx = px + ix * pw - w / 2
            cy = py + (1 - iy) * ph - h / 2
            self.rect(s, cx, cy, w, h, BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.5)
            self.text(s, cx, cy, w, h, [[(name, 12, True, "#FFFFFF")]],
                      align="center", anchor="middle")
        self._source(s, source)
        return s

    # ---------------------------------------------------------------- ⑬ 関係図(ラベル付き矢印)
    def flow(self, title, nodes, arrows=None, lead=None, band=None, source=None):
        """ノード(箱)をラベル付き矢印でつなぐ関係図。直列の役割分担などに使う。

        nodes : [(見出し, 小見出し|None, [項目...]), ...] 2〜3個
        arrows: 矢印の上に載せるラベル ["課題を引き渡す", ...](len(nodes)-1 個)
        band  : 下に敷く共通基盤バーの文言(省略可)
        """
        s = self._slide(LAYOUT_BODY, title=title)
        y0 = self._lead(s, lead) - 0.10
        n = len(nodes)
        gap = 2.05 if arrows else 1.45   # ラベル付き矢印はギャップを広げて収める
        w = (self.CW - gap * (n - 1)) / n
        band_h = 0.60
        bot = self.BODY_BOT - (band_h + 0.30 if band else 0.10)
        for i, (head, sub, items) in enumerate(nodes):
            x = self.MX + i * (w + gap)
            self.rect(s, x, y0, w, 0.55, BLUE)
            self.text(s, x, y0, w, 0.55, [[(head, 16, True, "#FFFFFF")]],
                      align="center", anchor="middle")
            yy = y0 + 0.55
            self.rect(s, x, yy, w, bot - yy, BG)  # 本文はヘッダに密着(1枚のカードに見せる)
            iy = yy + 0.16
            if sub:
                self._tag(s, x + 0.26, iy, sub)
                iy += 0.46
            self.text(s, x + 0.26, iy, w - 0.52, bot - iy - 0.15,
                      self.nest_paras(items, size=12), space_after=8)
            if i < n - 1:
                ax = x + w + 0.10
                aw = gap - 0.20
                ay = (y0 + bot) / 2
                tw = min(0.46, aw)
                self.tri_right(s, ax + (aw - tw) / 2, ay - 0.30, tw, 0.60, CYAN)
                if arrows and i < len(arrows) and arrows[i]:
                    # ラベルはギャップ内に収める(2行まで折返し可。箱への重なり防止)
                    self.text(s, ax - 0.12, ay - 1.10, aw + 0.24, 0.82,
                              [[(arrows[i], 10, True, BLUE_TX)]],
                              align="center", anchor="bottom")
        if band:
            self.rect(s, self.MX, bot + 0.22, self.CW, band_h, LTBLUE)
            self.text(s, self.MX + 0.2, bot + 0.22, self.CW - 0.4, band_h,
                      [[(band, 13, True, BLUE_TX)]], align="center", anchor="middle")
        self._source(s, source)
        return s

    # ---------------------------------------------------------------- ⑨ ロードマップ
    def timeline(self, title, phases, lead=None):
        s = self._slide(LAYOUT_BODY, title=title)
        y0 = self._lead(s, lead)
        n = len(phases)
        colors = RAMP[-n:] if n <= len(RAMP) else [RAMP[min(i, len(RAMP) - 1)] for i in range(n)]
        overlap = 0.24
        w = (self.CW + overlap * (n - 1)) / n
        ch_y = y0 + 0.18
        for i, (period, head, descs) in enumerate(phases):
            x = self.MX + i * (w - overlap)
            self.rect(s, x, ch_y, w, 0.88, colors[i],
                      shape=(MSO_SHAPE.CHEVRON if i > 0 else MSO_SHAPE.PENTAGON), adj=0.28)
            self.text(s, x + (0.30 if i > 0 else 0.12), ch_y, w - 0.42, 0.88,
                      [[(period, 15, True, "#FFFFFF")]], align="center", anchor="middle")
            bx = self.MX + i * (w - overlap) + (0.10 if i > 0 else 0)
            bw = w - overlap - 0.10
            by = ch_y + 1.10
            self.text(s, bx + 0.06, by, bw, 0.45, [[(head, 14, True, BLUE_TX)]])
            self.rect(s, bx + 0.06, by + 0.44, min(1.2, bw - 0.2), 0.035, CYAN)
            self.text(s, bx + 0.06, by + 0.62, bw - 0.10, self.BODY_BOT - by - 0.6,
                      self.nest_paras(descs, size=11), space_after=6)
        return s

    # ---------------------------------------------------------------- ⑩ メッセージ / 結び
    def message(self, text, sub=None, contact=None):
        s = self._slide(LAYOUT_LAST)
        size = 28 if len(text) <= 24 else 23
        paras = [[(text, size, True, NAVY)]]
        if sub:
            paras.append([(sub, 15, False, INK2)])
        self.text(s, 0.10 * self.W, 0.217 * self.H, 0.80 * self.W, 1.5,
                  paras, align="center", space_after=12)
        if contact:
            self.text(s, 0.10 * self.W, 0.787 * self.H, 0.80 * self.W, 1.0,
                      [[(c, 12, False, INK2)] for c in contact],
                      align="center", space_after=4)
        return s

    # ---------------------------------------------------------------- 保存
    def save(self, path):
        self.prs.save(path)
        return path
